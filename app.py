import os, hashlib, time, threading, logging, re, json
from datetime import datetime
from urllib.request import urlopen, Request
from html.parser import HTMLParser
import psycopg2, psycopg2.extras
from pptx import Presentation
import io, base64 as b64mod
from flask import Flask, jsonify, request
from flask_cors import CORS
from apscheduler.schedulers.background import BackgroundScheduler

logging.basicConfig(level=logging.INFO, format='%(asctime)s %(levelname)s %(message)s')
log = logging.getLogger(__name__)
app = Flask(__name__)
CORS(app)

DATABASE_URL = os.environ.get('DATABASE_URL', '')
ANTHROPIC_API_KEY = os.environ.get('ANTHROPIC_API_KEY', '')

# ── API ───────────────────────────────────────────────────────────────────────
def api_reorder_sources():
    """Save drag-drop order for sources"""
    data=request.get_json(force=True)
    orders=data.get('orders',[])  # [{url, sort_order, cat, region}]
    if not orders:
        return jsonify({'ok':True})
    conn=get_db(); cur=conn.cursor()
    for item in orders:
        cur.execute("""INSERT INTO source_order(url,cat,region,sort_order)
            VALUES(%s,%s,%s,%s)
            ON CONFLICT(url) DO UPDATE SET cat=EXCLUDED.cat,region=EXCLUDED.region,sort_order=EXCLUDED.sort_order""",
            (item.get('url'),item.get('cat',''),item.get('region',''),item.get('sort_order',0)))
    conn.commit(); conn.close()
    return jsonify({'ok':True})

@app.route('/consultant')
def consultant():
    return CONSULTANT_PAGE, 200, {"Content-Type": "text/html; charset=utf-8"}


LANDING_PAGE = """<!DOCTYPE html>
<html lang="fr">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>SubstanCiel</title>
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Syne:wght@400;600;700;800&family=DM+Sans:ital,wght@0,300;0,400;0,500;1,400&display=swap" rel="stylesheet">
<style>
*, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }

:root {
  --accent:  #1a3c2e;
  --accent2: #1f4a38;
  --accent3: #2a5c46;
  --lime:    #c8e84e;
  --lime2:   #b0d035;
  --lime-bg: rgba(200,232,78,0.10);
  --bg:      #f2f4f0;
  --surface: #ffffff;
  --text:    #111a14;
  --text2:   #3a4a3e;
  --muted:   #7a8e80;
  --border:  #e0e5d8;
  --shadow:  0 2px 8px rgba(26,60,46,0.08);
  --shadow-md: 0 6px 24px rgba(26,60,46,0.11);
  --shadow-lg: 0 16px 48px rgba(26,60,46,0.15);
}

html, body {
  height: 100%;
  font-family: 'DM Sans', system-ui, sans-serif;
  background: var(--accent);
  -webkit-font-smoothing: antialiased;
  overflow: hidden;
}

/* ── PAGE ── */
.page {
  height: 100vh;
  display: flex;
  flex-direction: column;
  position: relative;
  overflow: hidden;
}

/* ── NOISE TEXTURE ── */
.page::before {
  content: '';
  position: absolute; inset: 0;
  background-image: url("data:image/svg+xml,%3Csvg viewBox='0 0 256 256' xmlns='http://www.w3.org/2000/svg'%3E%3Cfilter id='noise'%3E%3CfeTurbulence type='fractalNoise' baseFrequency='0.9' numOctaves='4' stitchTiles='stitch'/%3E%3C/filter%3E%3Crect width='100%25' height='100%25' filter='url(%23noise)' opacity='0.03'/%3E%3C/svg%3E");
  pointer-events: none; z-index: 0;
}

/* ── GLOW BLOBS ── */
.blob {
  position: absolute;
  border-radius: 50%;
  filter: blur(100px);
  opacity: 0.12;
  pointer-events: none;
  z-index: 0;
}
.blob-1 { width: 600px; height: 600px; background: var(--lime); top: -200px; left: 50%; transform: translateX(-50%); animation: float 14s ease-in-out infinite alternate; }
.blob-2 { width: 300px; height: 300px; background: #5adf7a; bottom: 0; left: -60px; animation: float 10s ease-in-out infinite alternate-reverse; }
.blob-3 { width: 200px; height: 200px; background: var(--lime2); bottom: 80px; right: 60px; animation: float 8s ease-in-out infinite alternate; }

@keyframes float { 0% { transform: translateY(0) scale(1); } 100% { transform: translateY(20px) scale(1.05); } }
.blob-1 { animation: floatCenter 14s ease-in-out infinite alternate; }
@keyframes floatCenter { 0% { transform: translateX(-50%) translateY(0); } 100% { transform: translateX(-50%) translateY(24px); } }

/* ── HEADER ── */
header {
  position: relative; z-index: 10;
  display: flex; align-items: center;
  padding: 24px 48px;
  border-bottom: 1px solid rgba(255,255,255,0.06);
}

.logo {
  display: flex; align-items: center; gap: 10px;
}
.logo-mark {
  width: 34px; height: 34px;
  background: var(--lime);
  border-radius: 8px;
  display: flex; align-items: center; justify-content: center;
}
.logo-mark svg { width: 17px; height: 17px; }
.logo-name {
  font-family: 'Syne', sans-serif;
  font-weight: 800; font-size: 17px;
  color: #fff; letter-spacing: -0.3px;
}
.logo-name span { color: var(--lime); }

.header-pill {
  margin-left: auto;
  display: inline-flex; align-items: center; gap: 6px;
  font-size: 10.5px; font-weight: 600;
  color: rgba(200,232,78,0.65);
  letter-spacing: 0.09em; text-transform: uppercase;
  border: 1px solid rgba(200,232,78,0.15);
  padding: 5px 12px; border-radius: 100px;
}
.pulse {
  width: 6px; height: 6px;
  background: var(--lime2); border-radius: 50%;
  animation: pulse 2.2s ease-in-out infinite;
}
@keyframes pulse { 0%,100%{opacity:1;transform:scale(1)} 50%{opacity:.35;transform:scale(.7)} }

/* ── HERO ── */
.hero {
  flex: 1;
  display: flex; flex-direction: column;
  align-items: center; justify-content: center;
  text-align: center;
  padding: 0 24px 24px;
  position: relative; z-index: 10;
  gap: 0;
}

.hero-eyebrow {
  font-size: 11px; font-weight: 600;
  color: rgba(255,255,255,0.35);
  letter-spacing: 0.12em; text-transform: uppercase;
  margin-bottom: 20px;
}

h1 {
  font-family: 'Syne', sans-serif;
  font-size: clamp(42px, 5.5vw, 68px);
  font-weight: 800;
  color: #fff;
  line-height: 1.0;
  letter-spacing: -2px;
  margin-bottom: 20px;
}
h1 .lime { color: var(--lime); }
h1 .dim  { color: rgba(255,255,255,0.25); font-weight: 700; }

.hero-desc {
  font-size: 15px; font-weight: 300;
  color: rgba(255,255,255,0.4);
  line-height: 1.65;
  max-width: 420px;
  margin-bottom: 44px;
}

/* ── CARDS ── */
.cards {
  display: flex; gap: 14px;
  width: 100%; max-width: 640px;
}

.card {
  flex: 1;
  text-decoration: none;
  border-radius: 18px;
  padding: 24px 26px;
  display: flex; flex-direction: column;
  transition: transform 0.22s cubic-bezier(0.16,1,0.3,1), box-shadow 0.22s;
  position: relative; overflow: hidden;
}
.card:hover { transform: translateY(-4px); }

.card-primary {
  background: var(--lime);
}
.card-primary:hover {
  box-shadow: 0 18px 48px rgba(200,232,78,0.25);
}

.card-secondary {
  background: rgba(255,255,255,0.06);
  border: 1px solid rgba(255,255,255,0.10);
  backdrop-filter: blur(12px);
}
.card-secondary:hover {
  background: rgba(255,255,255,0.10);
  border-color: rgba(255,255,255,0.18);
  box-shadow: 0 18px 48px rgba(0,0,0,0.2);
}

.card-header {
  display: flex; align-items: center;
  justify-content: space-between;
  margin-bottom: 14px;
}
.card-icon {
  width: 38px; height: 38px;
  border-radius: 10px;
  display: flex; align-items: center; justify-content: center;
  font-size: 17px;
}
.card-primary .card-icon { background: rgba(26,60,46,0.12); }
.card-secondary .card-icon { background: rgba(255,255,255,0.08); }

.card-arrow {
  font-size: 20px;
  transition: transform 0.2s;
}
.card-primary .card-arrow { color: var(--accent); }
.card-secondary .card-arrow { color: rgba(255,255,255,0.4); }
.card:hover .card-arrow { transform: translate(3px,-3px); }

.card-title {
  font-family: 'Syne', sans-serif;
  font-size: 16px; font-weight: 800;
  letter-spacing: -0.3px;
  margin-bottom: 6px;
}
.card-primary .card-title { color: var(--accent); }
.card-secondary .card-title { color: #fff; }

.card-desc {
  font-size: 12px; line-height: 1.55;
}
.card-primary .card-desc { color: rgba(26,60,46,0.6); }
.card-secondary .card-desc { color: rgba(255,255,255,0.38); }

.card-tags {
  display: flex; flex-wrap: wrap; gap: 5px;
  margin-top: 16px;
}
.tag {
  font-size: 9.5px; font-weight: 700;
  padding: 3px 8px; border-radius: 100px;
  letter-spacing: 0.05em;
  text-transform: uppercase;
}
.card-primary .tag { background: rgba(26,60,46,0.1); color: var(--accent); }
.card-secondary .tag { background: rgba(255,255,255,0.08); color: rgba(255,255,255,0.45); border: 1px solid rgba(255,255,255,0.08); }

/* ── FOOTER ── */
footer {
  position: relative; z-index: 10;
  text-align: center;
  padding: 16px;
  font-size: 10.5px; color: rgba(255,255,255,0.18);
  letter-spacing: 0.06em;
}

@media (max-width: 600px) {
  header { padding: 18px 24px; }
  h1 { font-size: 36px; letter-spacing: -1px; }
  .cards { flex-direction: column; max-width: 380px; }
  html, body { overflow: auto; }
  .page { height: auto; min-height: 100vh; }
}
</style>
</head>
<body>
<div class="page">

  <div class="blob blob-1"></div>
  <div class="blob blob-2"></div>
  <div class="blob blob-3"></div>

  <!-- HEADER -->
  <header>
    <div class="logo">
      <div class="logo-mark">
        <svg viewBox="0 0 24 24" fill="none">
          <path d="M12 2L3 7v5c0 4.97 3.8 9.63 9 10.93C17.2 21.63 21 16.97 21 12V7L12 2z" fill="#1a3c2e"/>
          <path d="M8.5 12l2.5 2.5 4.5-5" stroke="#c8e84e" stroke-width="2.2" stroke-linecap="round" stroke-linejoin="round"/>
        </svg>
      </div>
      <div class="logo-name">Substan<span>Ciel</span></div>
    </div>
    <div class="header-pill">
      <span class="pulse"></span>
      Veille active
    </div>
  </header>

  <!-- HERO -->
  <div class="hero">
    <p class="hero-eyebrow">Financement public · Intelligence artificielle</p>

    <h1>
      Les bons financements<br>
      <span class="lime">au bon moment</span><br>
      <span class="dim">pour vos clients</span>
    </h1>

    <p class="hero-desc">
      Agrégation de subventions et appels à projets nationaux et régionaux — qualifiés et structurés par IA pour les consultants en financement.
    </p>

    <!-- CARTES -->
    <div class="cards">

      <a href="/consultant" class="card card-primary">
        <div class="card-header">
          <div class="card-icon">🔭</div>
          <span class="card-arrow">↗</span>
        </div>
        <div class="card-title">Espace Veille</div>
        <div class="card-desc">Parcourez, filtrez et qualifiez les dispositifs de financement en temps réel.</div>
        <div class="card-tags">
          <span class="tag">Curation IA</span>
          <span class="tag">70+ sources</span>
          <span class="tag">Multi-régions</span>
        </div>
      </a>

      <a href="/consultant" class="card card-secondary">
        <div class="card-header">
          <div class="card-icon">📋</div>
          <span class="card-arrow">↗</span>
        </div>
        <div class="card-title">Espace Collecte</div>
        <div class="card-desc">Collectez et exportez les fiches. Pré-veille 360° et journal par client.</div>
        <div class="card-tags">
          <span class="tag">Pré-veille 360°</span>
          <span class="tag">Export PPTX</span>
          <span class="tag">Journal</span>
        </div>
      </a>

    </div>
  </div>

  <footer>SubstanCiel · Outil interne de veille subventions</footer>

</div>
</body>
</html>
"""


@app.route('/')
def index():
    return LANDING_PAGE, 200, {"Content-Type": "text/html; charset=utf-8"}

@app.route('/app')
def app_page():
    return HTML_PAGE, 200, {"Content-Type": "text/html; charset=utf-8"}

@app.route('/api/ping')
def ping():
    return 'pong', 200



# ══════════════════════════════════════════════════════════════════
# AUTO-TAG AGENT
# ══════════════════════════════════════════════════════════════════

@app.route('/api/collect', methods=['POST'])
def collect_dispositif():
    """Fetch a URL, send to Claude, return structured grid."""
    data = request.get_json()
    url = data.get('url','')
    title = data.get('title','')
    article_id = data.get('id')
    if not url:
        return jsonify({'error':'URL required'}),400
    if not ANTHROPIC_API_KEY:
        return jsonify({'error':'ANTHROPIC_API_KEY not configured'}),500

    page_text = ''
    pdf_url = data.get('pdf_url', '')
    source_used = 'page'

    if article_id and not pdf_url:
        try:
            conn_tmp = get_db(); cur_tmp = conn_tmp.cursor()
            cur_tmp.execute("SELECT pdf_url FROM articles WHERE id=%s", (article_id,))
            row_tmp = cur_tmp.fetchone()
            if row_tmp and row_tmp['pdf_url']:
                pdf_url = row_tmp['pdf_url']
            cur_tmp.close(); conn_tmp.close()
        except Exception:
            pass

    if not pdf_url:
        try:
            pdf_url = _scrape_pdf_url(url)
        except Exception:
            pass

    # Priorite 1 : CDC PDF (timeout 12s)
    if pdf_url and pdf_url.lower().split('?')[0].endswith(('.pdf','.doc','.docx')):
        try:
            req_cdc = Request(pdf_url, headers={'User-Agent':'Mozilla/5.0'})
            with urlopen(req_cdc, timeout=12) as resp_cdc:
                raw_cdc = resp_cdc.read(150000)
            try:
                from io import BytesIO
                from pdfminer.high_level import extract_text as pdf_extract
                page_text = pdf_extract(BytesIO(raw_cdc))[:6000]
                source_used = 'cdc_pdf'
            except Exception:
                page_text = raw_cdc.decode('utf-8', errors='ignore')[:6000]
                source_used = 'cdc_raw'
        except Exception as e:
            log.warning(f"CDC fetch error {pdf_url}: {e}")

    # Priorite 2 : page HTML (timeout 10s) — extraction intelligente du contenu utile
    if not page_text:
        try:
            req_html = Request(url, headers={
                'User-Agent':'Mozilla/5.0 (Windows NT 10.0; Win64; x64)',
                'Accept-Language':'fr-FR,fr;q=0.9',
            })
            with urlopen(req_html, timeout=10) as resp_html:
                raw_html = resp_html.read(200000).decode('utf-8', errors='ignore')

            # Supprimer scripts, styles, nav, footer (bruit)
            NOISE_PAT = re.compile('<(script|style|nav|header|footer|aside)[^>]*>.*?</(script|style|nav|header|footer|aside)>', re.IGNORECASE|re.DOTALL)
            clean = NOISE_PAT.sub(' ', raw_html)

            # Essayer d'extraire la zone de contenu principal
            CONTENT_PAT = re.compile('<(main|article|section|div)[^>]*(content|main|article|body|dispositif|fiche|detail|description)[^>]*>(.*?)</(main|article|section|div)>', re.IGNORECASE|re.DOTALL)
            main_match = CONTENT_PAT.search(clean)
            if main_match:
                zone = main_match.group(3)
            else:
                zone = clean  # fallback : tout le HTML nettoyé

            # Strip tags restants
            text = re.sub(r'<[^>]+>', ' ', zone)
            text = re.sub(r'\s+', ' ', text).strip()

            # Garder 8000 chars — sauter les 500 premiers (souvent menu/breadcrumb)
            if len(text) > 500:
                text = text[500:]
            page_text = text[:8000]

        except Exception as e:
            log.warning(f"Fetch error {url}: {e}")
            page_text = f"Titre : {title}\nURL : {url}\n(Contenu non accessible)"

    # Call Claude Haiku (timeout 25s)
    try:
        cdc_mention = f"\nCahier des charges : {pdf_url}" if pdf_url else ""
        user_content = f"Analyse ce dispositif et remplis la grille.{cdc_mention}\n\nTitre : {title}\nURL : {url}\n[Source : {source_used}]\n\nContenu :\n{page_text}"
        payload = json.dumps({
            "model": "claude-haiku-4-5-20251001",
            "max_tokens": 2000,
            "system": COLLECT_PROMPT,
            "messages": [{"role":"user","content":user_content}]
        }, ensure_ascii=False).encode('utf-8')
        req = Request("https://api.anthropic.com/v1/messages", data=payload, headers={
            "Content-Type":"application/json; charset=utf-8",
            "x-api-key":ANTHROPIC_API_KEY,
            "anthropic-version":"2023-06-01"
        }, method="POST")
        with urlopen(req, timeout=30) as resp:
            raw_resp = resp.read()
            claude_data = json.loads(raw_resp)
        if claude_data.get('type') == 'error':
            raise Exception(f"Anthropic API error: {claude_data.get('error',{}).get('message','unknown')}")
        text = claude_data["content"][0]["text"].strip()
        m = re.search(r'\{[\s\S]*\}', text)
        result = json.loads(m.group() if m else text)
        result['source_url'] = url
        result['article_id'] = article_id
        if pdf_url:
            result['cdc_url'] = pdf_url
        return jsonify(result)
    except Exception as e:
        import traceback
        log.error(f"Collect Claude error: {e}\n{traceback.format_exc()}")
        return jsonify({'error': str(e)}),500




# ═══════════════════════════════════════════════════════════════════════════════
# PACKAGES
# ═══════════════════════════════════════════════════════════════════════════════

@app.route('/api/packages', methods=['GET'])
def get_packages():
    conn = get_db(); cur = conn.cursor()
    cur.execute("""
        SELECT p.id, p.name, p.created_at,
               COUNT(d.id) as nb
        FROM packages p
        LEFT JOIN dispositifs d ON d.package_id = p.id
        GROUP BY p.id ORDER BY p.created_at DESC
    """)
    rows = cur.fetchall(); cur.close(); conn.close()
    result = []
    for r in rows:
        result.append({'id': r['id'], 'name': r['name'],
                       'created_at': r['created_at'].isoformat() if r['created_at'] else '',
                       'nb': r['nb']})
    return jsonify(result)

@app.route('/api/packages', methods=['POST'])
def create_package():
    data = request.get_json()
    name = data.get('name','').strip()
    if not name:
        return jsonify({'error': 'Nom requis'}), 400
    conn = get_db(); cur = conn.cursor()
    cur.execute("INSERT INTO packages (name) VALUES (%s) RETURNING id", (name,))
    pkg_id = cur.fetchone()['id']
    conn.commit(); cur.close(); conn.close()
    return jsonify({'id': pkg_id, 'name': name})

@app.route('/api/packages/<int:pid>', methods=['DELETE'])
def delete_package(pid):
    conn = get_db(); cur = conn.cursor()
    cur.execute("DELETE FROM packages WHERE id=%s", (pid,))
    conn.commit(); cur.close(); conn.close()
    return jsonify({'status': 'deleted'})

@app.route('/api/packages/<int:pid>/dispositifs', methods=['GET'])
def get_package_dispositifs(pid):
    conn = get_db(); cur = conn.cursor()
    cur.execute("SELECT * FROM dispositifs WHERE package_id=%s ORDER BY id ASC", (pid,))
    rows = cur.fetchall(); cur.close(); conn.close()
    result = []
    for r in rows:
        d = dict(r)
        if d.get('collected_at'): d['collected_at'] = d['collected_at'].isoformat()
        result.append(d)
    return jsonify(result)



@app.route('/api/packages/merge', methods=['POST'])
def merge_packages():
    """Merge two packages into a new one."""
    data = request.get_json()
    pkg_a = data.get('pkg_a')  # source package (current)
    pkg_b = data.get('pkg_b')  # target package to merge with
    new_name = data.get('name', '').strip()
    if not pkg_a or not pkg_b or not new_name:
        return jsonify({'error': 'Paramètres manquants'}), 400
    if pkg_a == pkg_b:
        return jsonify({'error': 'Impossible de fusionner un package avec lui-même'}), 400

    conn = get_db(); cur = conn.cursor()
    # Create new package
    cur.execute("INSERT INTO packages (name) VALUES (%s) RETURNING id", (new_name,))
    new_id = cur.fetchone()['id']
    # Move all dispositifs from A and B into new package (deduplicate by source_url)
    cur.execute("""
        INSERT INTO dispositifs (guichet_financeur, guichet_instructeur, titre, nature,
            beneficiaire, type_depot, date_fermeture, objectif, types_depenses,
            operations_eligibles, depenses_eligibles, criteres_eligibilite,
            depenses_ineligibles, montants_taux, thematiques, territoire,
            points_vigilance, contact, programme_europeen, source_url, cdc_url, package_id)
        SELECT DISTINCT ON (COALESCE(source_url, gen_random_uuid()::text))
            guichet_financeur, guichet_instructeur, titre, nature,
            beneficiaire, type_depot, date_fermeture, objectif, types_depenses,
            operations_eligibles, depenses_eligibles, criteres_eligibilite,
            depenses_ineligibles, montants_taux, thematiques, territoire,
            points_vigilance, contact, programme_europeen, source_url, cdc_url, %s
        FROM dispositifs
        WHERE package_id IN (%s, %s)
        ORDER BY id ASC, COALESCE(source_url, gen_random_uuid()::text), collected_at DESC
    """, (new_id, pkg_a, pkg_b))
    # Delete source packages (dispositifs cascade to SET NULL, already moved)
    cur.execute("DELETE FROM dispositifs WHERE package_id IN (%s, %s)", (pkg_a, pkg_b))
    cur.execute("DELETE FROM packages WHERE id IN (%s, %s)", (pkg_a, pkg_b))
    conn.commit(); cur.close(); conn.close()
    return jsonify({'status': 'merged', 'new_id': new_id, 'name': new_name})


@app.route('/api/packages/<int:pid>/logs', methods=['GET'])
def get_package_logs(pid):
    """Return error logs from batch jobs linked to this package."""
    conn = get_db(); cur = conn.cursor()
    cur.execute("""
        SELECT job_id, created_at, total, done, results
        FROM batch_jobs
        WHERE pkg_id = %s
        ORDER BY created_at DESC
        LIMIT 10
    """, (pid,))
    rows = cur.fetchall(); cur.close(); conn.close()
    logs = []
    for r in rows:
        results = r['results'] or []
        if isinstance(results, str):
            import json as _json
            results = _json.loads(results)
        errors = [x for x in results if x.get('status') == 'error']
        logs.append({
            'job_id': r['job_id'],
            'created_at': r['created_at'].isoformat() if r['created_at'] else '',
            'total': r['total'],
            'done': r['done'],
            'errors': errors
        })
    return jsonify(logs)

@app.route('/api/packages/<int:pid>/export-cdc', methods=['GET'])
def export_package_cdc(pid):
    """Download all CDC documents for a package as a ZIP."""
    conn = get_db(); cur = conn.cursor()
    cur.execute("SELECT name FROM packages WHERE id=%s", (pid,))
    pkg = cur.fetchone()
    if not pkg:
        return jsonify({'error': 'Package introuvable'}), 404
    cur.execute("SELECT titre, source_url, cdc_url FROM dispositifs WHERE package_id=%s AND cdc_url IS NOT NULL AND cdc_url != ''", (pid,))
    rows = cur.fetchall(); cur.close(); conn.close()
    if not rows:
        return jsonify({'error': 'Aucun CDC trouvé dans ce package'}), 404

    import zipfile, io as _io
    from urllib.request import Request as _Req, urlopen as _open
    buf = _io.BytesIO()
    added = 0
    with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zf:
        for r in rows:
            titre = (r['titre'] or 'dispositif').replace('/', '-').replace('\\', '-')[:50]
            cdc_url = r['cdc_url']
            try:
                ext = cdc_url.split('?')[0].rsplit('.', 1)[-1].lower()
                if ext not in ('pdf', 'doc', 'docx', 'odt'):
                    ext = 'pdf'
                req = _Req(cdc_url, headers={'User-Agent': 'Mozilla/5.0'})
                with _open(req, timeout=15) as resp:
                    data = resp.read(10_000_000)  # 10 Mo max
                safe_name = f"{added+1:02d}_{titre}.{ext}"
                zf.writestr(safe_name, data)
                added += 1
            except Exception as e:
                log.warning(f"CDC download error {cdc_url}: {e}")
                continue

    if added == 0:
        return jsonify({'error': 'Impossible de télécharger les CDCs'}), 500

    buf.seek(0)
    from flask import send_file
    safe_pkg = pkg['name'].replace(' ', '_').replace('/', '-')[:40]
    return send_file(buf, mimetype='application/zip',
                     as_attachment=True,
                     download_name=f"CDCs_{safe_pkg}.zip")

@app.route('/api/packages/<int:pid>/export-pptx', methods=['GET'])
def export_package_pptx(pid):
    try:
     return _export_package_pptx_inner(pid)
    except Exception as e:
        import traceback
        log.error(f"export_package_pptx fatal: {e}\n{traceback.format_exc()}")
        return jsonify({'error': str(e)}), 500

def _export_package_pptx_inner(pid):
    conn = get_db(); cur = conn.cursor()
    cur.execute("SELECT name FROM packages WHERE id=%s", (pid,))
    pkg = cur.fetchone()
    if not pkg:
        return jsonify({'error': 'Package introuvable'}), 404
    cur.execute("SELECT * FROM dispositifs WHERE package_id=%s ORDER BY id ASC", (pid,))
    rows = cur.fetchall(); cur.close(); conn.close()
    if not rows:
        return jsonify({'error': 'Package vide'}), 400

    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    import io, base64 as b64mod

    import copy

    def _merge_slide_into_prs(src_slide, dst_prs):
        """Copy a slide including its images into dst_prs."""
        layout = dst_prs.slide_layouts[5]
        new_slide = dst_prs.slides.add_slide(layout)

        # Copy image parts using get_or_add_image_part (returns (ImagePart, rId))
        rId_map = {}
        for rel in src_slide.part.rels.values():
            if 'image' in rel.reltype:
                try:
                    src_img = rel.target_part
                    _img_part, new_rId = new_slide.part.get_or_add_image_part(
                        io.BytesIO(src_img.blob)
                    )
                    rId_map[rel.rId] = new_rId
                except Exception as e:
                    log.warning(f"Image copy error rId={rel.rId}: {e}")

        # Copy spTree with rId remapping for images
        src_sp_tree = src_slide.shapes._spTree
        dst_sp_tree = new_slide.shapes._spTree

        # Clear destination placeholders (keep first 2 mandatory group nodes)
        while len(dst_sp_tree) > 2:
            dst_sp_tree.remove(dst_sp_tree[-1])

        # Deep-copy each child, remapping r:embed / r:link attributes
        for child in list(src_sp_tree)[2:]:
            el = copy.deepcopy(child)
            for node in el.iter():
                for attr in list(node.attrib.keys()):
                    if attr.endswith('}embed') or attr.endswith('}link'):
                        old_rId = node.attrib[attr]
                        if old_rId in rId_map:
                            node.attrib[attr] = rId_map[old_rId]
            dst_sp_tree.append(el)

    # Generate all individual PPTX bytes
    all_pptx = []
    for r in rows:
        data = dict(r)
        if data.get('collected_at'): data['collected_at'] = data['collected_at'].isoformat()
        try:
            pptx_bytes = generate_dispositif_pptx(data)
            if pptx_bytes:
                all_pptx.append(pptx_bytes)
            else:
                log.warning(f"Package PPTX generate returned None for id={data.get('id')}")
        except Exception as e:
            import traceback
            log.error(f"Package PPTX generate error id={data.get('id')} titre={data.get('titre','?')}: {e}\n{traceback.format_exc()}")
            continue

    if not all_pptx:
        return jsonify({'error': 'Aucune slide generee'}), 500

    # Use first pptx as base, merge all others into it
    base_prs = Presentation(io.BytesIO(all_pptx[0]))

    for pptx_bytes in all_pptx[1:]:
        try:
            src_prs = Presentation(io.BytesIO(pptx_bytes))
            for slide in src_prs.slides:
                _merge_slide_into_prs(slide, base_prs)
        except Exception as e:
            log.warning(f"Package PPTX merge error: {e}")
            continue

    combined_prs = base_prs

    try:
        buf = io.BytesIO()
        combined_prs.save(buf)
        pptx_data = buf.getvalue()
    except Exception as e:
        log.error(f"Package PPTX save error: {e}")
        import traceback; log.error(traceback.format_exc())
        return jsonify({'error': f'Erreur sauvegarde PPTX: {str(e)}'}), 500

    safe_name = (pkg['name'] or 'package').replace(' ', '_').replace('/', '-')[:40]
    from flask import make_response
    resp = make_response(pptx_data)
    resp.headers['Content-Type'] = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    resp.headers['Content-Disposition'] = f'attachment; filename="Package_{safe_name}.pptx"'
    resp.headers['Content-Length'] = len(pptx_data)
    return resp


# ═══════════════════════════════════════════════════════════════════════════════
# BATCH COLLECT (Excel upload)
# ═══════════════════════════════════════════════════════════════════════════════


# ── Batch collect state (DB-backed, multi-worker safe) ───────────────────────
def _job_update(job_id, done=None, result=None, status=None):
    """Atomically update a batch job in DB."""
    try:
        conn = get_db(); cur = conn.cursor()
        if result is not None:
            cur.execute(
                "UPDATE batch_jobs SET done=done+1, results=results||%s::jsonb WHERE job_id=%s",
                (json.dumps([result]), job_id)
            )
        if status:
            cur.execute("UPDATE batch_jobs SET status=%s WHERE job_id=%s", (status, job_id))
        conn.commit(); cur.close(); conn.close()
    except Exception as e:
        log.error(f"job_update error: {e}")

def _job_get(job_id):
    try:
        conn = get_db(); cur = conn.cursor()
        cur.execute("SELECT * FROM batch_jobs WHERE job_id=%s", (job_id,))
        row = cur.fetchone(); cur.close(); conn.close()
        if not row: return None
        return dict(row)
    except Exception:
        return None

def _job_create(job_id, total, pkg_id, pkg_name):
    try:
        conn = get_db(); cur = conn.cursor()
        cur.execute(
            "INSERT INTO batch_jobs (job_id, status, total, done, pkg_id, pkg_name, results) VALUES (%s,'running',%s,0,%s,%s,'[]')",
            (job_id, total, pkg_id, pkg_name)
        )
        conn.commit(); cur.close(); conn.close()
    except Exception as e:
        log.error(f"job_create error: {e}")



@app.route('/api/collect-text', methods=['POST'])
def collect_text():
    """Analyze raw pasted text content with Claude."""
    if not ANTHROPIC_API_KEY:
        return jsonify({'error': 'ANTHROPIC_API_KEY not configured'}), 500
    data = request.get_json()
    text = (data.get('text') or '').strip()
    source_url = (data.get('source_url') or '').strip()
    if not text:
        return jsonify({'error': 'Contenu vide'}), 400
    try:
        url_mention = f"\nURL source : {source_url}" if source_url else ""
        user_content = f"Analyse ce contenu et remplis la grille.{url_mention}\n[Source : scrape_manuel]\n\nContenu :\n{text[:8000]}"
        payload = json.dumps({
            "model": "claude-haiku-4-5-20251001",
            "max_tokens": 2000,
            "system": COLLECT_PROMPT,
            "messages": [{"role": "user", "content": user_content}]
        }).encode()
        req = Request("https://api.anthropic.com/v1/messages", data=payload, headers={
            "Content-Type": "application/json",
            "x-api-key": ANTHROPIC_API_KEY,
            "anthropic-version": "2023-06-01"
        }, method="POST")
        with urlopen(req, timeout=30) as resp:
            claude_data = json.loads(resp.read())
        txt = claude_data["content"][0]["text"].strip()
        m = re.search(r'\{[\s\S]*\}', txt)
        result = json.loads(m.group() if m else txt)
        result['source_url'] = source_url or ''
        return jsonify(result)
    except Exception as e:
        log.error(f"collect_text error: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/collect-cdc', methods=['POST'])
def collect_cdc():
    """Analyze an uploaded CDC file (PDF/Word) directly with Claude."""
    if not ANTHROPIC_API_KEY:
        return jsonify({'error': 'ANTHROPIC_API_KEY not configured'}), 500

    file = request.files.get('file')
    if not file:
        return jsonify({'error': 'Fichier manquant'}), 400

    source_url = request.form.get('source_url', '').strip()
    filename = file.filename.lower()

    page_text = ''
    source_used = 'cdc_upload'

    try:
        raw = file.read(200000)
        if filename.endswith('.pdf'):
            try:
                from io import BytesIO
                from pdfminer.high_level import extract_text as pdf_extract
                page_text = pdf_extract(BytesIO(raw))[:8000]
            except Exception:
                page_text = raw.decode('utf-8', errors='ignore')[:8000]
        else:
            # Word doc - try docx
            try:
                import zipfile, io as _io
                with zipfile.ZipFile(_io.BytesIO(raw)) as z:
                    if 'word/document.xml' in z.namelist():
                        xml = z.read('word/document.xml').decode('utf-8', errors='ignore')
                        import re as _re
                        page_text = _re.sub(r'<[^>]+>', ' ', xml)
                        page_text = _re.sub(r'\s+', ' ', page_text).strip()[:8000]
            except Exception:
                page_text = raw.decode('utf-8', errors='ignore')[:8000]
    except Exception as e:
        return jsonify({'error': f'Lecture fichier impossible : {e}'}), 400

    if not page_text.strip():
        return jsonify({'error': 'Impossible d extraire le texte du document'}), 400

    try:
        url_mention = f"\nURL source : {source_url}" if source_url else ""
        user_content = f"Analyse ce cahier des charges et remplis la grille.{url_mention}\n[Source : {source_used}]\n\nContenu :\n{page_text}"
        payload = json.dumps({
            "model": "claude-haiku-4-5-20251001",
            "max_tokens": 2000,
            "system": COLLECT_PROMPT,
            "messages": [{"role": "user", "content": user_content}]
        }).encode()
        req = Request("https://api.anthropic.com/v1/messages", data=payload, headers={
            "Content-Type": "application/json",
            "x-api-key": ANTHROPIC_API_KEY,
            "anthropic-version": "2023-06-01"
        }, method="POST")
        with urlopen(req, timeout=30) as resp:
            claude_data = json.loads(resp.read())
        text = claude_data["content"][0]["text"].strip()
        m = re.search(r'\{[\s\S]*\}', text)
        result = json.loads(m.group() if m else text)
        result['source_url'] = source_url or ''
        result['cdc_uploaded'] = True
        return jsonify(result)
    except Exception as e:
        log.error(f"CDC collect error: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/collect-batch', methods=['POST'])
def collect_batch():
    """Start async batch collect. Returns job_id immediately."""
    try:
        import openpyxl
    except ImportError:
        return jsonify({'error': 'openpyxl non installe'}), 500

    file = request.files.get('file')
    if not file:
        return jsonify({'error': 'Fichier manquant'}), 400

    package_name = request.form.get('package_name', '').strip()
    create_pkg = request.form.get('create_package', 'false') == 'true' and bool(package_name)

    try:
        import io as _io
        wb = openpyxl.load_workbook(_io.BytesIO(file.read()), read_only=True, data_only=True)
        ws = wb.worksheets[0]
        urls = []
        for row in ws.iter_rows(min_row=1, max_row=31, min_col=1, max_col=1, values_only=True):
            val = row[0]
            if val and isinstance(val, str) and val.strip().startswith('http'):
                urls.append(val.strip())
        wb.close()
    except Exception as e:
        return jsonify({'error': f'Lecture Excel impossible : {e}'}), 400

    if not urls:
        return jsonify({'error': 'Aucune URL trouvee en colonne A'}), 400
    urls = urls[:30]

    # Create package if requested
    pkg_id = None
    if create_pkg:
        conn = get_db(); cur = conn.cursor()
        cur.execute("INSERT INTO packages (name) VALUES (%s) RETURNING id", (package_name,))
        pkg_id = cur.fetchone()['id']
        conn.commit(); cur.close(); conn.close()

    import uuid
    job_id = str(uuid.uuid4())[:8]
    _job_create(job_id, len(urls), pkg_id, package_name)

    # Run in background thread
    def run_job():
        fields = ['guichet_financeur','guichet_instructeur','titre','nature','beneficiaire',
                  'type_depot','date_fermeture','objectif','types_depenses','operations_eligibles',
                  'depenses_eligibles','criteres_eligibilite','depenses_ineligibles','montants_taux',
                  'thematiques','territoire','points_vigilance','contact','programme_europeen','source_url','cdc_url']
        for idx, url in enumerate(urls):
            result = {'url': url, 'index': idx, 'status': 'error', 'titre': '', 'error': ''}
            try:
                page_text = ''
                pdf_url = None
                source_used = 'page'
                try:
                    pdf_url = _scrape_pdf_url(url)
                except Exception:
                    pass
                if pdf_url and pdf_url.lower().split('?')[0].endswith(('.pdf','.doc','.docx')):
                    try:
                        req_cdc = Request(pdf_url, headers={'User-Agent':'Mozilla/5.0'})
                        with urlopen(req_cdc, timeout=12) as resp_cdc:
                            raw_cdc = resp_cdc.read(150000)
                        try:
                            from io import BytesIO
                            from pdfminer.high_level import extract_text as pdf_extract
                            page_text = pdf_extract(BytesIO(raw_cdc))[:6000]
                            source_used = 'cdc_pdf'
                        except Exception:
                            page_text = raw_cdc.decode('utf-8', errors='ignore')[:6000]
                            source_used = 'cdc_raw'
                    except Exception as e:
                        log.warning(f"Batch CDC error {pdf_url}: {e}")
                if not page_text:
                    try:
                        req_html = Request(url, headers={'User-Agent':'Mozilla/5.0 (Windows NT 10.0; Win64; x64)'})
                        with urlopen(req_html, timeout=10) as resp_html:
                            raw_html = resp_html.read(200000).decode('utf-8', errors='ignore')
                        NOISE_PAT = re.compile('<(script|style|nav|header|footer|aside)[^>]*>.*?</(script|style|nav|header|footer|aside)>', re.IGNORECASE|re.DOTALL)
                        clean = NOISE_PAT.sub(' ', raw_html)
                        text = re.sub(r'<[^>]+>', ' ', clean)
                        text = re.sub(r'\s+', ' ', text).strip()
                        page_text = text[500:8500] if len(text) > 500 else text[:8000]
                    except Exception:
                        page_text = f"URL: {url}"
                cdc_mention = f"\nCahier des charges : {pdf_url}" if pdf_url else ""
                user_content = f"Analyse ce dispositif et remplis la grille.{cdc_mention}\nURL : {url}\n[Source : {source_used}]\n\nContenu :\n{page_text}"
                payload = json.dumps({
                    "model": "claude-haiku-4-5-20251001",
                    "max_tokens": 2000,
                    "system": COLLECT_PROMPT,
                    "messages": [{"role":"user","content":user_content}]
                }, ensure_ascii=False).encode('utf-8')
                req = Request("https://api.anthropic.com/v1/messages", data=payload, headers={
                    "Content-Type":"application/json",
                    "x-api-key": ANTHROPIC_API_KEY,
                    "anthropic-version":"2023-06-01"
                }, method="POST")
                with urlopen(req, timeout=30) as resp:
                    claude_data = json.loads(resp.read())
                text_resp = claude_data["content"][0]["text"].strip()
                m = re.search(r'\{[\s\S]*\}', text_resp)
                disp = json.loads(m.group() if m else text_resp)
                disp['source_url'] = url
                if pdf_url: disp['cdc_url'] = pdf_url
                conn2 = get_db(); cur2 = conn2.cursor()
                # Only block duplicate within the same package (or globally if no package)
                if pkg_id:
                    cur2.execute("SELECT id FROM dispositifs WHERE source_url=%s AND package_id=%s", (url, pkg_id))
                else:
                    cur2.execute("SELECT id FROM dispositifs WHERE source_url=%s AND package_id IS NULL", (url,))
                existing = cur2.fetchone()
                if existing:
                    result['status'] = 'duplicate'
                    result['titre'] = disp.get('titre', url)
                else:
                    cols = ','.join(fields)
                    placeholders = ','.join(['%s']*len(fields))
                    vals = [disp.get(f,'') for f in fields]
                    if pkg_id:
                        cur2.execute(f"INSERT INTO dispositifs ({cols}, package_id) VALUES ({placeholders}, %s) RETURNING id", vals + [pkg_id])
                    else:
                        cur2.execute(f"INSERT INTO dispositifs ({cols}) VALUES ({placeholders}) RETURNING id", vals)
                    conn2.commit()
                    result['status'] = 'saved'
                    result['titre'] = disp.get('titre', url)
                cur2.close(); conn2.close()
            except Exception as e:
                result['error'] = str(e)[:120]
                log.error(f"Batch error {url}: {e}")
            _job_update(job_id, result=result)
        _job_update(job_id, status='done')

    t = threading.Thread(target=run_job, daemon=True)
    t.start()
    return jsonify({'job_id': job_id, 'total': len(urls), 'pkg_id': pkg_id, 'pkg_name': package_name})


@app.route('/api/collect-batch/<job_id>', methods=['GET'])
def collect_batch_status(job_id):
    """Poll batch collect job status."""
    job = _job_get(job_id)
    if not job:
        return jsonify({'error': 'Job introuvable'}), 404
    # Convert DB row to expected format
    results = job.get('results') or []
    if isinstance(results, str):
        results = json.loads(results)
    return jsonify({
        'status': job['status'],
        'total': job['total'],
        'done': job['done'],
        'pkg_id': job['pkg_id'],
        'pkg_name': job['pkg_name'],
        'results': results
    })


@app.route('/api/dispositifs', methods=['GET'])
def get_dispositifs():
    conn = get_db(); cur = conn.cursor()
    # Deduplicate ONLY on non-empty source_url, keep all manual entries, order by id ASC
    cur.execute("""
        SELECT DISTINCT ON (CASE WHEN source_url IS NOT NULL AND source_url != '' THEN source_url ELSE id::text END) *
        FROM dispositifs
        ORDER BY CASE WHEN source_url IS NOT NULL AND source_url != '' THEN source_url ELSE id::text END, id ASC
    """)
    rows = cur.fetchall(); cur.close(); conn.close()
    result = []
    for r in rows:
        d = dict(r)
        if d.get('collected_at'): d['collected_at'] = d['collected_at'].isoformat()
        result.append(d)
    return jsonify(result)


@app.route('/api/dispositifs/<int:did>', methods=['DELETE'])
def delete_dispositif(did):
    conn = get_db(); cur = conn.cursor()
    cur.execute("DELETE FROM dispositifs WHERE id=%s", (did,))
    conn.commit(); cur.close(); conn.close()
    return jsonify({'status': 'deleted'})

@app.route('/api/dispositifs', methods=['POST'])
def save_dispositif():
    data = request.get_json()
    fields = ['guichet_financeur','guichet_instructeur','titre','nature','beneficiaire',
              'type_depot','date_fermeture','objectif','types_depenses','operations_eligibles',
              'depenses_eligibles','criteres_eligibilite','depenses_ineligibles','montants_taux',
              'thematiques','territoire','points_vigilance','contact','programme_europeen','source_url']
    conn = get_db(); cur = conn.cursor()
    cols = ','.join(fields)
    placeholders = ','.join(['%s']*len(fields))
    vals = [data.get(f,'') for f in fields]
    src_url = data.get('source_url','')
    if src_url:
        cur.execute("SELECT id FROM dispositifs WHERE source_url=%s", (src_url,))
        if cur.fetchone():
            cur.close(); conn.close()
            return jsonify({'status':'duplicate','message':'Déjà dans la base'}), 200
    cur.execute(f"INSERT INTO dispositifs ({cols}) VALUES ({placeholders})", vals)
    conn.commit(); cur.close(); conn.close()
    return jsonify({'status':'saved'})


# ═══════════════════════════════════════════════════════════════════════════════
# PDF / CAHIERS DES CHARGES
# ═══════════════════════════════════════════════════════════════════════════════

CDC_DOC_KEYWORDS = [
    'cahier', 'cahier-des-charges', 'reglement', 'regl', 'appel-a-projets',
    'appel_a_projets', 'notice', 'dossier', 'formulaire', 'guide', 'annexe',
    'modalites', 'candidature', 'depot', 'programme', 'cdc', 'specifications'
]
CDC_DOC_EXTENSIONS = ('.pdf', '.doc', '.docx', '.odt', '.xls', '.xlsx')

def _make_absolute(href, page_url):
    """Convert relative href to absolute URL."""
    from urllib.parse import urlparse, urljoin
    if href.startswith('http'):
        return href
    return urljoin(page_url, href)

def _scrape_pdf_url(page_url):
    """Visit a page and find a CDC document link (PDF/Word/image). Returns URL or None."""
    try:
        req = Request(page_url, headers={
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36',
            'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8',
        })
        with urlopen(req, timeout=8) as resp:
            raw = resp.read(400000).decode('utf-8', errors='replace')

        # Extraire tous les liens <a href="...">texte</a>
        links = re.findall(r'<a[^>]+href=["\'\s]?([^"\'\s>]+)["\'\s>][^>]*>(.*?)</a>',
                           raw, re.IGNORECASE | re.DOTALL)

        candidates_url_kw  = []  # href contient extension + mot-clé
        candidates_url_ext = []  # href contient juste extension
        candidates_txt_kw  = []  # texte du lien contient mot-clé

        for href, text in links:
            href = href.strip()
            if not href or href.startswith('#') or href.startswith('mailto'):
                continue
            lower_href = href.lower().split('?')[0]
            text_clean = re.sub(r'<[^>]+>', ' ', text).strip().lower()
            has_ext = any(lower_href.endswith(ext) for ext in CDC_DOC_EXTENSIONS)
            has_kw_url = any(kw in lower_href for kw in CDC_DOC_KEYWORDS)
            has_kw_txt = any(kw in text_clean for kw in CDC_DOC_KEYWORDS)

            abs_href = _make_absolute(href, page_url)
            if not abs_href.startswith('http'):
                continue

            if has_ext and has_kw_url:
                candidates_url_kw.append(abs_href)
            elif has_ext:
                candidates_url_ext.append(abs_href)
            elif has_kw_txt and has_ext:
                # Texte CDC + extension document = candidat valide
                candidates_txt_kw.append(abs_href)
            # Sinon : lien HTML avec texte CDC = ignoré (trop de bruit)

        # Retourne le meilleur candidat par priorité
        # Seuls les liens avec vraie extension document sont retenus
        for pool in [candidates_url_kw, candidates_url_ext, candidates_txt_kw]:
            if pool:
                return pool[0]

    except Exception as e:
        log.warning(f"CDC scrape failed for {page_url}: {e}")
    return None

def _scrape_pdf_url_ai(page_url):
    """Use Claude to find the PDF/CDC link on a page. Returns URL or None."""
    try:
        req = Request(page_url, headers={'User-Agent': 'Mozilla/5.0'})
        with urlopen(req, timeout=10) as resp:
            raw = resp.read(100000).decode('utf-8', errors='replace')
        # Strip tags for Claude
        clean = re.sub(r'<[^>]+>', ' ', raw)
        clean = re.sub(r'\s+', ' ', clean)[:6000]

        payload = json.dumps({
            "model": "claude-haiku-4-5-20251001",
            "max_tokens": 300,
            "messages": [{
                "role": "user",
                "content": f"""Analyse cette page web et trouve l'URL du cahier des charges, règlement, ou document PDF principal (appel à projets, dossier de candidature, notice, etc.).

URL de la page : {page_url}

Contenu de la page (extrait) :
{clean}

Réponds UNIQUEMENT avec l'URL complète du PDF si tu en trouves un. Si tu n'en trouves pas, réponds exactement : AUCUN"""
            }]
        }).encode()

        api_req = Request(
            'https://api.anthropic.com/v1/messages',
            data=payload,
            headers={
                'Content-Type': 'application/json',
                'x-api-key': ANTHROPIC_API_KEY,
                'anthropic-version': '2023-06-01'
            }
        )
        with urlopen(api_req, timeout=30) as resp:
            result = json.loads(resp.read())
        text = result['content'][0]['text'].strip()
        if text and text != 'AUCUN' and text.startswith('http'):
            return text
    except Exception as e:
        log.warning(f"AI PDF search failed: {e}")
    return None

@app.route('/api/articles/fetch-pdf', methods=['POST'])
def fetch_pdf_single():
    """Scraping pour 1 article avec debug détaillé."""
    data = request.json or {}
    article_id = data.get('article_id')
    if not article_id:
        return jsonify({'error': 'article_id required'}), 400
    try:
        conn = get_db(); cur = conn.cursor()
        cur.execute("SELECT url, title FROM articles WHERE id=%s", (article_id,))
        row = cur.fetchone()
        if not row:
            cur.close(); conn.close()
            return jsonify({'error': 'not found'}), 404
        page_url = row['url']
        doc_url = None
        debug_info = {'page_url': page_url, 'links_found': 0, 'error': None}
        try:
            req = Request(page_url, headers={
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36',
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8',
            })
            with urlopen(req, timeout=8) as resp:
                raw = resp.read(400000).decode('utf-8', errors='replace')
            links = re.findall(r'<a[^>]+href=["\'\s]?([^"\'\s>]+)["\'\s>][^>]*>(.*?)</a>',
                               raw, re.IGNORECASE | re.DOTALL)
            debug_info['links_found'] = len(links)
            # Collect all candidates for debug
            all_ext_links = []
            for href, text in links:
                href = href.strip()
                if not href or href.startswith('#') or href.startswith('mailto'):
                    continue
                lower_href = href.lower().split('?')[0]
                if any(lower_href.endswith(ext) for ext in CDC_DOC_EXTENSIONS):
                    abs_href = _make_absolute(href, page_url)
                    text_clean = re.sub(r'<[^>]+>', ' ', text).strip()[:60]
                    all_ext_links.append({'href': abs_href, 'text': text_clean})
            debug_info['ext_links'] = all_ext_links[:10]
            doc_url = _scrape_pdf_url(page_url)
        except Exception as e:
            debug_info['error'] = str(e)
            log.error(f"CDC scrape error #{article_id}: {e}")
        cur.execute("UPDATE articles SET pdf_url=%s WHERE id=%s", (doc_url, article_id))
        conn.commit(); cur.close(); conn.close()
        log.info(f"CDC #{article_id}: {doc_url} | debug={debug_info}")
        return jsonify({'article_id': article_id, 'pdf_url': doc_url, 'doc_url': doc_url,
                        'title': row['title'], 'debug': debug_info})
    except Exception as e:
        log.error(f"CDC route error: {e}")
        return jsonify({'error': str(e), 'pdf_url': None, 'doc_url': None}), 200

# ── CDC scan job state ──────────────────────────────────────────────────────
_cdc_job = {'status': 'idle', 'done': 0, 'total': 0, 'results': [], 'error': None}
_cdc_lock = threading.Lock()

def _run_cdc_scan_bg(article_ids, use_ai=False):
    """Background thread : scan CDC sans bloquer Gunicorn."""
    from concurrent.futures import ThreadPoolExecutor, as_completed as fut_completed
    global _cdc_job
    with _cdc_lock:
        _cdc_job = {'status': 'running', 'done': 0, 'total': len(article_ids), 'results': [], 'error': None}

    conn = get_db(); cur = conn.cursor()
    try:
        articles = []
        for aid in article_ids:
            cur.execute("SELECT id, url, title, pdf_url FROM articles WHERE id=%s", (aid,))
            row = cur.fetchone()
            if row:
                articles.append(dict(row))

        def scan_one(art):
            try:
                if use_ai:
                    doc_url = art['pdf_url'] or _scrape_pdf_url_ai(art['url'])
                else:
                    doc_url = _scrape_pdf_url(art['url'])
                return {'article_id': art['id'], 'doc_url': doc_url, 'title': art['title'], 'source': 'ai' if use_ai else 'scan'}
            except Exception as e:
                return {'article_id': art['id'], 'doc_url': None, 'title': art.get('title',''), 'source': 'error'}

        results = []
        with ThreadPoolExecutor(max_workers=6) as ex:
            futures = {ex.submit(scan_one, a): a for a in articles}
            for fut in fut_completed(futures):
                r = fut.result()
                results.append(r)
                # Save to DB immediately
                try:
                    cur.execute("UPDATE articles SET pdf_url=%s WHERE id=%s", (r['doc_url'], r['article_id']))
                    conn.commit()
                except Exception:
                    conn.rollback()
                with _cdc_lock:
                    _cdc_job['done'] += 1
                    _cdc_job['results'].append(r)

        with _cdc_lock:
            _cdc_job['status'] = 'done'
    except Exception as e:
        with _cdc_lock:
            _cdc_job['status'] = 'error'
            _cdc_job['error'] = str(e)
        log.error(f"CDC scan error: {e}")
    finally:
        cur.close(); conn.close()

@app.route('/api/articles/fetch-pdf-batch', methods=['POST'])
def fetch_pdf_batch():
    """Lance un scan CDC en arrière-plan et retourne immédiatement."""
    global _cdc_job
    with _cdc_lock:
        if _cdc_job['status'] == 'running':
            return jsonify({'status': 'already_running', 'done': _cdc_job['done'], 'total': _cdc_job['total']}), 200

    data = request.json or {}
    ids = data.get('article_ids', [])
    if not ids:
        return jsonify({'error': 'article_ids required'}), 400

    ids = ids[:200]  # max 200 articles
    t = threading.Thread(target=_run_cdc_scan_bg, args=(ids, False), daemon=True)
    t.start()
    return jsonify({'status': 'started', 'total': len(ids)})

@app.route('/api/articles/fetch-pdf-status', methods=['GET'])
def fetch_pdf_status():
    """Polling : retourne l'état du scan CDC en cours."""
    with _cdc_lock:
        job = dict(_cdc_job)
    return jsonify(job)

@app.route('/api/articles/fetch-pdf-ai', methods=['POST'])
def fetch_pdf_ai():
    """Lance un scan IA CDC en arrière-plan."""
    global _cdc_job
    if not ANTHROPIC_API_KEY:
        return jsonify({'error': 'API key not configured'}), 500
    with _cdc_lock:
        if _cdc_job['status'] == 'running':
            return jsonify({'status': 'already_running'}), 200
    data = request.json or {}
    ids = data.get('article_ids', [])
    if not ids:
        return jsonify({'error': 'article_ids required'}), 400
    ids = ids[:30]  # max 30 pour l'IA (coût)
    t = threading.Thread(target=_run_cdc_scan_bg, args=(ids, True), daemon=True)
    t.start()
    return jsonify({'status': 'started', 'total': len(ids)})

init_db()
